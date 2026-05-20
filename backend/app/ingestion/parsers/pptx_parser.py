import io
import logging
from typing import Dict, Any, List, Tuple

logger = logging.getLogger(__name__)

try:
    import pptx
except ImportError:
    pptx = None


def extract_pptx_content(file_bytes: bytes, filename: str) -> Tuple[str, List[Dict[str, Any]], Dict[str, Any]]:
    """
    Parses a PPTX document slide by slide, extracting text, slide titles, and speaker notes.
    """
    full_text_list = []
    slides_meta = []
    global_meta = {"filename": filename}

    if not pptx:
        logger.error("python-pptx is not installed.")
        return file_bytes.decode("utf-8", errors="ignore"), [], global_meta

    try:
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_title = ""
            slide_texts = []
            
            # 1. Try to get slide title
            if slide.shapes.title:
                slide_title = slide.shapes.title.text.strip()
            
            # 2. Get text from all shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            slide_texts.append(p_text)
            
            # 3. Get speaker notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            # Compile text for this slide
            slide_header = f"--- SLIDE {slide_num}: {slide_title or 'Untitled'} ---"
            slide_body = "\n".join(slide_texts)
            slide_notes_block = f"\n[Speaker Notes]\n{notes_text}" if notes_text else ""
            
            combined_slide_text = f"{slide_header}\n{slide_body}{slide_notes_block}"
            full_text_list.append(combined_slide_text)
            
            slides_meta.append({
                "slide_number": slide_num,
                "section_title": slide_title or f"Slide {slide_num}",
                "char_count": len(combined_slide_text)
            })

    except Exception as e:
        logger.error(f"Error parsing PPTX {filename}: {e}")
        return file_bytes.decode("utf-8", errors="ignore"), [], global_meta

    full_text = "\n\n".join(full_text_list)
    return full_text, slides_meta, global_meta
